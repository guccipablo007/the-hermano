#!/usr/bin/env python3
import argparse, json, re, subprocess, zipfile
from pathlib import Path
from xml.sax.saxutils import escape
MODEL_CALL=Path('/root/.hermes/scripts/hermes_model_call.py'); VERIFY=Path('/root/.hermes/scripts/hermes_verify.py'); DELIVER=Path('/root/.hermes/scripts/hermes_telegram_deliver.py')
def run(cmd, timeout=180): return subprocess.run(cmd, shell=True, text=True, capture_output=True, timeout=timeout)
def call_basic_worker(title, description, expect_text):
    prompt=f"""You are the Basic Research and Summary Worker.
Create clean document content in Markdown format.
Title: {title}
Task: {description}
Requirements:
- Start with this exact title as an H1 heading: {title}
- Include the expected text exactly somewhere: {expect_text or title}
- Use clear headings, short paragraphs, and useful bullet points.
- Do not include markdown code fences.
- Do not include fake file paths.
- Do not claim you created a file.
- Return only the document content.
"""
    proc=run(f"{MODEL_CALL} --route basic --prompt {json.dumps(prompt)} --max-tokens 1600 --json", timeout=240)
    if proc.returncode!=0: raise RuntimeError('MODEL_CALL_FAILED\nSTDOUT='+proc.stdout+'\nSTDERR='+proc.stderr)
    data=json.loads(proc.stdout); content=(data.get('content') or '').strip()
    if content.startswith('```'):
        lines=content.splitlines()
        if lines and lines[0].startswith('```'): lines=lines[1:]
        if lines and lines[-1].strip().startswith('```'): lines=lines[:-1]
        content='\n'.join(lines).strip()
    if not content: raise RuntimeError('EMPTY_MODEL_CONTENT')
    if expect_text and expect_text not in content: content=f"# {title}\n\n{expect_text}\n\n"+content
    return content,data
def blocks(md): return [b.strip() for b in re.split(r'\n\s*\n', md.strip()) if b.strip()]
def save_md(md,out): out.write_text(md+'\n')
def save_docx(md,out):
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc=Document(); sec=doc.sections[0]; sec.top_margin=int(.6*914400); sec.bottom_margin=int(.6*914400); sec.left_margin=int(.7*914400); sec.right_margin=int(.7*914400); doc.styles['Normal'].font.name='Arial'; doc.styles['Normal'].font.size=Pt(11)
    for block in blocks(md):
        lines=block.splitlines()
        if lines[0].startswith('# '): p=doc.add_heading(lines[0][2:].strip(),0); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
        elif lines[0].startswith('## '): doc.add_heading(lines[0][3:].strip(),1)
        elif lines[0].startswith('### '): doc.add_heading(lines[0][4:].strip(),2)
        elif all(l.strip().startswith(('- ','* ')) for l in lines if l.strip()):
            for l in lines: doc.add_paragraph(l.strip()[2:].strip(), style='List Bullet')
        else: doc.add_paragraph(' '.join(l.strip() for l in lines))
    doc.save(out)
    if not zipfile.is_zipfile(out): raise RuntimeError('DOCX_ZIP_VALIDATION_FAILED')
def save_pdf(md,out):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
    from reportlab.lib import colors
    doc=SimpleDocTemplate(str(out), pagesize=A4, rightMargin=.65*inch, leftMargin=.65*inch, topMargin=.65*inch, bottomMargin=.65*inch); styles=getSampleStyleSheet()
    styles.add(ParagraphStyle(name='HermesTitle', parent=styles['Title'], fontSize=22, leading=27, alignment=1, textColor=colors.HexColor('#0f172a'), spaceAfter=18)); styles.add(ParagraphStyle(name='HermesHeading', parent=styles['Heading2'], fontSize=15, leading=20, textColor=colors.HexColor('#1e293b'), spaceBefore=14, spaceAfter=8)); styles.add(ParagraphStyle(name='HermesBody', parent=styles['BodyText'], fontSize=10.5, leading=15, textColor=colors.HexColor('#111827'), spaceAfter=8))
    story=[]
    for block in blocks(md):
        lines=block.splitlines()
        if lines[0].startswith('# '): story.append(Paragraph(escape(lines[0][2:].strip()), styles['HermesTitle'])); story.append(Spacer(1,.08*inch))
        elif lines[0].startswith('## '): story.append(Paragraph(escape(lines[0][3:].strip()), styles['HermesHeading']))
        elif lines[0].startswith('### '): story.append(Paragraph(escape(lines[0][4:].strip()), styles['HermesHeading']))
        elif all(l.strip().startswith(('- ','* ')) for l in lines if l.strip()): story.append(ListFlowable([ListItem(Paragraph(escape(l.strip()[2:].strip()), styles['HermesBody'])) for l in lines], bulletType='bullet', leftIndent=18)); story.append(Spacer(1,.06*inch))
        else: story.append(Paragraph(escape(' '.join(l.strip() for l in lines)), styles['HermesBody']))
    doc.build(story)
    if not out.read_bytes().startswith(b'%PDF'): raise RuntimeError('PDF_HEADER_VALIDATION_FAILED')

def save_pptx(md,out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    def add_title_slide(title_text, subtitle='Generated by Hermes'):
        slide=prs.slides.add_slide(prs.slide_layouts[6]); fill=slide.background.fill; fill.solid(); fill.fore_color.rgb=RGBColor(15,23,42)
        box=slide.shapes.add_textbox(Inches(.8), Inches(1.8), Inches(11.8), Inches(1.2)); tf=box.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=title_text; p.alignment=PP_ALIGN.CENTER; r=p.runs[0]; r.font.size=Pt(44); r.font.bold=True; r.font.color.rgb=RGBColor(255,255,255)
        sub=slide.shapes.add_textbox(Inches(1.5), Inches(3.15), Inches(10.3), Inches(.8)); stf=sub.text_frame; stf.clear(); sp=stf.paragraphs[0]; sp.text=subtitle; sp.alignment=PP_ALIGN.CENTER; sr=sp.runs[0]; sr.font.size=Pt(20); sr.font.color.rgb=RGBColor(203,213,225)
    def add_content_slide(heading, items):
        slide=prs.slides.add_slide(prs.slide_layouts[6]); title_box=slide.shapes.add_textbox(Inches(.65), Inches(.45), Inches(12), Inches(.85)); tf=title_box.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=heading; r=p.runs[0]; r.font.size=Pt(31); r.font.bold=True; r.font.color.rgb=RGBColor(15,23,42)
        body=slide.shapes.add_textbox(Inches(.95), Inches(1.55), Inches(11.2), Inches(5.3)); btf=body.text_frame; btf.clear(); btf.word_wrap=True; clean=[x.strip() for x in items if x.strip()] or [' ']
        for i,item in enumerate(clean[:7]):
            para=btf.paragraphs[0] if i==0 else btf.add_paragraph(); para.text=item; para.level=0; para.font.size=Pt(22); para.space_after=Pt(10)
    title='Presentation'
    for block in blocks(md):
        lines=[x.strip() for x in block.splitlines() if x.strip()]
        if lines and lines[0].startswith('# '): title=lines[0][2:].strip(); break
    add_title_slide(title)
    heading=None; items=[]
    def flush():
        nonlocal heading, items
        if heading:
            add_content_slide(heading, items); heading=None; items=[]
    for block in blocks(md):
        lines=[x.strip() for x in block.splitlines() if x.strip()]
        if not lines: continue
        first=lines[0]
        if first.startswith('# '): continue
        if first.startswith('## '):
            flush(); heading=first[3:].strip(); rest=lines[1:]
            for x in rest:
                if x.startswith(('- ','* ')): x=x[2:].strip()
                items.append(re.sub(r'\*\*(.*?)\*\*', r'\1', x))
        elif first.startswith('### '):
            if not heading: heading=first[4:].strip()
            else: items.append(first[4:].strip())
        else:
            if not heading: heading='Overview'
            for x in lines:
                if x.startswith(('- ','* ')): x=x[2:].strip()
                items.append(re.sub(r'\*\*(.*?)\*\*', r'\1', x))
    flush()
    if len(prs.slides)<2: add_content_slide('Overview',['This presentation was generated by Hermes.'])
    prs.save(out)
    if not zipfile.is_zipfile(out): raise RuntimeError('PPTX_ZIP_VALIDATION_FAILED')

def verify(path, expect, fmt):
    if not path.exists() or path.stat().st_size<50: raise RuntimeError('OUTPUT_INVALID')
    if fmt=='pdf' and not path.read_bytes().startswith(b'%PDF'): raise RuntimeError('PDF_HEADER_MISSING')
    if fmt=='docx' and not zipfile.is_zipfile(path): raise RuntimeError('DOCX_ZIP_INVALID')
    if fmt=='pptx' and not zipfile.is_zipfile(path): raise RuntimeError('PPTX_ZIP_INVALID')
    if fmt=='md' and expect and expect not in path.read_text(errors='ignore'): raise RuntimeError('EXPECTED_TEXT_MISSING')
    proc=run(f"{VERIFY} file --path {json.dumps(str(path))}")
    if proc.returncode!=0: raise RuntimeError('HERMES_VERIFY_FAILED')
def deliver(path, fmt):
    proc=run(f"{DELIVER} --file {json.dumps(str(path))} --caption {json.dumps('Your Majesty, here is your verified '+fmt.upper()+' document: '+path.name)} --preview-link no", timeout=180)
    if proc.returncode!=0: raise RuntimeError('TELEGRAM_DELIVERY_FAILED\n'+proc.stdout+proc.stderr)
    return proc.stdout
def main():
    ap=argparse.ArgumentParser(); ap.add_argument('--format', choices=['md','docx','pdf','pptx'], required=True); ap.add_argument('--output', required=True); ap.add_argument('--title', required=True); ap.add_argument('--description', required=True); ap.add_argument('--content-file', default=''); ap.add_argument('--expect-text', default=''); ap.add_argument('--deliver-telegram', action='store_true'); ap.add_argument('--no-deliver-telegram', action='store_true'); a=ap.parse_args(); out=Path(a.output); out.parent.mkdir(parents=True, exist_ok=True)
    
    if a.content_file:
        content_path=Path(a.content_file)
        if not content_path.exists(): raise RuntimeError('CONTENT_FILE_NOT_FOUND')
        md=content_path.read_text(errors='replace').strip()
        if not md: raise RuntimeError('CONTENT_FILE_EMPTY')
        if not md.lstrip().startswith('#'):
            md='# '+a.title+'\n\n'+md
        data={'model':'direct-content-file'}
    else:
        md,data=call_basic_worker(a.title,a.description,a.expect_text)
    {'md':save_md,'docx':save_docx,'pdf':save_pdf,'pptx':save_pptx}[a.format](md,out)
    verify(out,a.expect_text,a.format); should=not a.no_deliver_telegram or a.deliver_telegram; delivery=''
    if should: delivery=deliver(out,a.format)
    print('DOCUMENT_BUILD=PASSED'); print('FORMAT='+a.format); print('ROUTE=basic'); print('WORKER_MODEL='+str(data.get('model'))); print('OUTPUT_PATH='+str(out)); print('EXPECT_TEXT='+(a.expect_text or '')); print('VERIFICATION=PASSED'); print('TELEGRAM_DELIVERY='+('PASSED' if should else 'SKIPPED'))
    if delivery: print('TELEGRAM_DELIVERY_OUTPUT_START'); print(delivery.strip()); print('TELEGRAM_DELIVERY_OUTPUT_END')
if __name__=='__main__': main()

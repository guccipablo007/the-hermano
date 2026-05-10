#!/usr/bin/env python3
import argparse
import json
import os
import math
import re
import subprocess
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

VERIFY = Path("/root/.hermes/scripts/hermes_verify.py")
MEDIA_VERIFY = Path("/root/.hermes/scripts/hermes_media_verify.py")
DELIVER = Path("/root/.hermes/scripts/hermes_telegram_deliver.py")

PALETTE = {
    "sky": (190, 225, 255), "blue": (80, 145, 255), "green": (90, 185, 125),
    "yellow": (255, 220, 95), "orange": (255, 160, 80), "red": (235, 85, 105),
    "pink": (255, 190, 220), "purple": (155, 120, 255), "brown": (120, 80, 45),
    "skin": (245, 190, 150), "dark": (35, 45, 65), "white": (255, 255, 255),
}

def run(cmd, timeout=180):
    return subprocess.run(cmd, shell=True, text=True, capture_output=True, timeout=timeout)

def q(s):
    return json.dumps(str(s))

def slug(text):
    return re.sub(r"[^a-z0-9]+", "_", (text or "slide").lower()).strip("_")[:60] or "slide"

def font(size, bold=False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for p in candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size=size)
    return ImageFont.load_default()

def base_canvas(title):
    img = Image.new("RGB", (960, 540), (248, 252, 255))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((22, 22, 938, 518), radius=36, fill=PALETTE["white"], outline=(180, 210, 255), width=7)
    d.text((44, 34), title[:26], font=font(46, True), fill=PALETTE["dark"])
    return img, d

def draw_person(d, x, y, shirt):
    d.ellipse((x+45, y+10, x+135, y+100), fill=PALETTE["skin"], outline=PALETTE["dark"], width=4)
    d.ellipse((x+70, y+45, x+82, y+57), fill=PALETTE["dark"])
    d.ellipse((x+100, y+45, x+112, y+57), fill=PALETTE["dark"])
    d.arc((x+73, y+60, x+112, y+86), 10, 170, fill=PALETTE["red"], width=4)
    d.rounded_rectangle((x+35, y+112, x+145, y+245), radius=30, fill=shirt, outline=PALETTE["dark"], width=4)

def draw_icon(path, title, image_key, topic=""):
    label = image_key or title
    tier = os.environ.get('HERMES_VISUAL_QUALITY_TIER','local_cartoon_scene')
    cmd = ['/usr/bin/python3','/root/.hermes/scripts/hermes_visual_asset_router.py','--label',str(label),'--topic',str(topic or title),'--output',str(path),'--quality-tier',tier]
    proc = subprocess.run(cmd, text=True, capture_output=True, timeout=180)
    if proc.returncode != 0:
        raise RuntimeError('VISUAL_ASSET_ROUTER_FAILED\n' + proc.stdout + proc.stderr)
    verify = subprocess.run(['/usr/bin/python3','/root/.hermes/scripts/hermes_visual_quality_verify.py','--image',str(path),'--min-tier',tier], text=True, capture_output=True, timeout=120)
    if verify.returncode != 0:
        raise RuntimeError('VISUAL_QUALITY_VERIFY_FAILED\n' + verify.stdout + verify.stderr)
def default_slides(topic, count):
    words = ["Wake Up", "Brush Teeth", "Wash Face", "Get Dressed", "Eat Breakfast", "Pack School Bag", "Go To School", "Learn New Things"]
    return [{"title": w, "sentence": f"I {w.lower()}.", "image": w.lower()} for w in words[:count]]

def load_slides(path, topic, count):
    if path:
        data = json.loads(Path(path).read_text())
        slides=[]
        for item in data:
            title = item.get("title") or item.get("word") or "Slide"
            slides.append({"title": title, "sentence": item.get("sentence") or f"I {title.lower()}.", "image": item.get("image") or title})
        return slides
    return default_slides(topic, count)

def build_pptx(output, title, topic, age, slides, min_media):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    img_dir = Path("/root/.hermes/file_outputs/pptx_picture_images") / output.stem
    img_dir.mkdir(parents=True, exist_ok=True)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 253, 248)
    tb = title_slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.5), Inches(1.1))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(46)
    r.font.bold = True
    r.font.color.rgb = RGBColor(15, 23, 42)
    sub = title_slide.shapes.add_textbox(Inches(1.6), Inches(3.0), Inches(10.2), Inches(0.6))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = f"Age: {age}"
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.runs[0]
    sr.font.size = Pt(24)
    sr.font.color.rgb = RGBColor(71, 85, 105)

    image_paths = []
    total = len(slides)
    for i, item in enumerate(slides, 1):
        img_path = img_dir / f"{i:02d}_{slug(item.get('title', 'slide'))}.png"
        draw_icon(img_path, item.get("image", item.get("title", "")), item.get("title", ""), topic)
        image_paths.append(img_path)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 253, 248)

        # image panel, left side
        slide.shapes.add_picture(str(img_path), Inches(0.65), Inches(1.25), width=Inches(6.2), height=Inches(4.65))

        # title
        title_box = slide.shapes.add_textbox(Inches(7.25), Inches(1.3), Inches(5.4), Inches(0.8))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = item.get("title", "")
        p.alignment = PP_ALIGN.LEFT
        r = p.runs[0]
        r.font.size = Pt(34)
        r.font.bold = True
        r.font.color.rgb = RGBColor(190, 24, 93)

        # sentence
        sentence_box = slide.shapes.add_textbox(Inches(7.25), Inches(2.35), Inches(5.35), Inches(1.2))
        stf = sentence_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = item.get("sentence", "")
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.runs[0]
        sr.font.size = Pt(25)
        sr.font.color.rgb = RGBColor(31, 41, 55)

        # speaking prompt
        prompt_box = slide.shapes.add_textbox(Inches(7.25), Inches(4.85), Inches(5.35), Inches(0.9))
        ptf = prompt_box.text_frame
        ptf.clear()
        pp = ptf.paragraphs[0]
        pp.text = "Can you say it?"
        pp.alignment = PP_ALIGN.LEFT
        pr = pp.runs[0]
        pr.font.size = Pt(24)
        pr.font.bold = True
        pr.font.color.rgb = RGBColor(37, 99, 235)

        # footer
        fb = slide.shapes.add_textbox(Inches(11.5), Inches(6.85), Inches(1.0), Inches(0.3))
        ftf = fb.text_frame
        ftf.clear()
        fp = ftf.paragraphs[0]
        fp.text = f"{i}/{total}"
        fp.alignment = PP_ALIGN.RIGHT
        fr = fp.runs[0]
        fr.font.size = Pt(12)
        fr.font.color.rgb = RGBColor(100, 116, 139)

    prs.save(output)
    if not zipfile.is_zipfile(output):
        raise RuntimeError("PPTX_INVALID")
    return image_paths


def main():
    ap=argparse.ArgumentParser(description='Build illustrated PPTX lessons with local PNG pictures.')
    ap.add_argument('--output', required=True); ap.add_argument('--title', required=True); ap.add_argument('--topic', default=''); ap.add_argument('--age', default='kids')
    ap.add_argument('--slide-count', type=int, default=8); ap.add_argument('--slides-json', default=''); ap.add_argument('--min-media', type=int, default=1)
    ap.add_argument('--quality-tier', default='local_cartoon_scene'); ap.add_argument('--deliver-telegram', action='store_true'); ap.add_argument('--no-deliver-telegram', action='store_true')
    args=ap.parse_args()
    os.environ['HERMES_VISUAL_QUALITY_TIER']=args.quality_tier
    slides=load_slides(args.slides_json, args.topic or args.title, args.slide_count)
    images_dir=Path('/root/.hermes/file_outputs/pptx_picture_images') / Path(args.output).stem
    image_paths=build_pptx(args.output, args.title, args.topic or args.title, args.age, slides, images_dir)
    vf=run(f"{VERIFY} file --path {q(args.output)}", timeout=120)
    if vf.returncode != 0:
        print('NOT VERIFIED'); print('REASON=FILE_VERIFY_FAILED'); print(vf.stdout); print(vf.stderr); return vf.returncode
    min_media=args.min_media or len(image_paths)
    mv=run(f"{MEDIA_VERIFY} --file {q(args.output)} --type pptx --min-media {int(min_media)}", timeout=120)
    if mv.returncode != 0:
        print('NOT VERIFIED'); print('REASON=MEDIA_VERIFY_FAILED'); print(mv.stdout); print(mv.stderr); return mv.returncode
    delivery='SKIPPED'
    if not args.no_deliver_telegram or args.deliver_telegram:
        dp=run(f"{DELIVER} --file {q(args.output)} --caption {q('Your Majesty, here is your verified illustrated PPTX lesson: ' + Path(args.output).name)} --preview-link no", timeout=180)
        if dp.returncode != 0:
            print('NOT VERIFIED'); print('REASON=TELEGRAM_DELIVERY_FAILED'); print(dp.stdout); print(dp.stderr); return dp.returncode
        delivery='PASSED'
    print('PPTX_PICTURE_BUILD=PASSED')
    print('ROUTE=pptx_picture_file')
    print('OUTPUT_PATH=' + str(args.output))
    print('SLIDE_COUNT=' + str(len(slides)))
    print('GENERATED_IMAGE_COUNT=' + str(len(image_paths)))
    print('IMAGES_DIR=' + str(images_dir))
    print('VERIFICATION=PASSED')
    print('MEDIA_VERIFY=PASSED')
    print('TELEGRAM_DELIVERY=' + delivery)
    print(mv.stdout.strip())
    return 0
if __name__ == '__main__':
    raise SystemExit(main())
